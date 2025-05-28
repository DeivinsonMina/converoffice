from flask import Flask, render_template, request, redirect, url_for, send_file
import os
from werkzeug.utils import secure_filename
from pdf2docx import Converter

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf'}

if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return redirect(request.url)
    file = request.files['file']
    if file.filename == '' or not allowed_file(file.filename):
        return redirect(request.url)
    
    format_selected = request.form.get('format')
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    if format_selected == 'docx':
        # Convert PDF to DOCX
        output_filename = filename.rsplit('.', 1)[0] + '.docx'
        output_filepath = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
        cv = Converter(filepath)
        cv.convert(output_filepath, start=0, end=None)
        cv.close()
    elif format_selected == 'xlsx':
        # Convert PDF to XLSX (simple text extraction)
        import pdfplumber
        from openpyxl import Workbook

        output_filename = filename.rsplit('.', 1)[0] + '.xlsx'
        output_filepath = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)

        wb = Workbook()
        ws = wb.active

        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        ws.append([line])

        wb.save(output_filepath)
    elif format_selected == 'pptx':
        # Convert PDF to PPTX (simple text extraction)
        import pdfplumber
        from pptx import Presentation

        output_filename = filename.rsplit('.', 1)[0] + '.pptx'
        output_filepath = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)

        prs = Presentation()
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = "Página del PDF"
                    slide.placeholders[1].text = text

        prs.save(output_filepath)
    else:
        return render_template('result.html', error="Formato no soportado.", filename=None)

    return render_template('result.html', filename=output_filename, error=None)

@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_file(os.path.join(app.config['UPLOAD_FOLDER'], filename), as_attachment=True)

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(debug=True, host='0.0.0.0', port=port)
